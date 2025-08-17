import io
import time
import zipfile
import re
from datetime import datetime, timezone
from typing import List, Dict

import streamlit as st
from PIL import Image, UnidentifiedImageError
import exifread

# PowerPoint
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# PDF processing
try:
    import fitz  # PyMuPDF
    PDF_SUPPORT = True
except ImportError:
    PDF_SUPPORT = False

# Camera detection constants
CAMERA_KEYS = [
    "EXIF LensModel", "EXIF FNumber", "EXIF FocalLength",
    "EXIF ExposureTime", "EXIF ISOSpeedRatings", "EXIF DateTimeOriginal",
    "Image Model", "Image Make", "GPS GPSLatitude", "GPS GPSLongitude"
]

# ---------- Helpers ----------
def timestamp_slug():
    return datetime.now(timezone.utc).strftime("%Y%m%dT%H%M%SZ")

def is_image(filename: str) -> bool:
    ext = (filename.rsplit(".", 1)[-1] if "." in filename else "").lower()
    return ext in {"png", "jpg", "jpeg", "gif", "webp", "bmp", "tiff", "mpo", "heic", "heif", "svg"}

def is_pdf(filename: str) -> bool:
    ext = (filename.rsplit(".", 1)[-1] if "." in filename else "").lower()
    return ext == "pdf"

def extract_file_metadata(file_bytes: bytes, filename: str) -> dict:
    """Extract camera detection from image files"""
    metadata = {
        "is_camera": False,
        "upload_time": datetime.now().isoformat()
    }
    
    if is_image(filename):
        try:
            # Use exifread to extract EXIF data
            tags = exifread.process_file(io.BytesIO(file_bytes), details=False)
            # Normalize dict to simple {key: str(value)}
            exif_data = {str(k): str(v) for k, v in tags.items()}
            
            # Check if camera EXIF fields are present
            present_cam_fields = [k for k in CAMERA_KEYS if k in exif_data]
            if present_cam_fields:
                metadata["is_camera"] = True
                st.info(f"📸 Camera detected in {filename}: {', '.join(present_cam_fields[:3])}...")
            else:
                st.info(f"📱 No camera EXIF found in {filename}")
                
        except Exception as e:
            st.warning(f"⚠️ Could not read EXIF from {filename}: {str(e)}")
    
    return metadata

def fit_image_on_slide(prs: Presentation, pil_img: Image.Image, has_title: bool = False):
    # Margins for portrait format
    max_w = prs.slide_width - Inches(1.0)  # 7.5" width available
    if has_title:
        # Leave space for smaller title at top
        max_h = prs.slide_height - Inches(2.5)  # 8.5" height available
        top_offset = Inches(1.5)  # Start below title area
    else:
        max_h = prs.slide_height - Inches(1.0)  # 10" height available
        top_offset = Inches(0.5)  # Center vertically
    
    iw, ih = pil_img.size
    # Convert EMUs to px scale factor by comparing to width/height as ints
    scale = min(float(max_w) / iw, float(max_h) / ih)
    disp_w = int(iw * scale)
    disp_h = int(ih * scale)
    left = int((prs.slide_width - disp_w) / 2)
    top = int(top_offset + (max_h - disp_h) / 2)
    return left, top, disp_w, disp_h

def process_image_for_powerpoint(image_bytes: bytes, filename: str) -> bytes:
    """Process image to ensure PowerPoint compatibility"""
    try:
        img = Image.open(io.BytesIO(image_bytes))
        
        # Convert unsupported formats to PNG
        if img.format not in ['JPEG', 'PNG', 'GIF', 'BMP', 'TIFF']:
            # Convert to PNG for PowerPoint compatibility
            img_bytes = io.BytesIO()
            img = img.convert('RGB')  # Convert to RGB if needed
            img.save(img_bytes, format='PNG')
            img_bytes.seek(0)
            st.info(f"🔄 Converted {filename} from {img.format} to PNG for PowerPoint compatibility")
            return img_bytes.read()
        
        # For supported formats, return original bytes
        return image_bytes
        
    except Exception as e:
        st.warning(f"⚠️ Could not process image '{filename}': {str(e)}")
        return image_bytes  # Return original bytes as fallback

def add_footer_to_slide(slide, batch_id: str, slide_width: int, slide_height: int):
    """Add batch ID footer to a slide"""
    if batch_id:
        # Add footer textbox at bottom
        footer_box = slide.shapes.add_textbox(Inches(0.5), slide_height - Inches(0.8), 
                                            slide_width - Inches(1.0), Inches(0.5))
        footer_frame = footer_box.text_frame
        footer_frame.clear()
        p = footer_frame.paragraphs[0]
        p.text = f"Batch: {batch_id}"
        p.font.size = Pt(10)
        p.font.color.rgb = RGBColor(128, 128, 128)  # Gray color
        p.alignment = PP_ALIGN.CENTER

def build_pptx_from_records(records: List[Dict], show_titles: bool = False, batch_id: str = "", presentation_title: str = "") -> bytes:
    """
    Creates a PPTX entirely in-memory:
      - image/* -> individual slides with image fitted
      - PDFs -> each page becomes a slide (if PyMuPDF available)
      - everything else -> one "Attached Files" summary slide listing names
    """
    prs = Presentation()
    
    # Set slide size to portrait/letter format (8.5" x 11")
    prs.slide_width = Inches(8.5)
    prs.slide_height = Inches(11.0)
    
    # Set presentation properties
    if presentation_title:
        prs.core_properties.title = presentation_title
        prs.core_properties.author = "U2P - Upload to Present"
        prs.core_properties.subject = f"Document Collection - {presentation_title}"
        prs.core_properties.keywords = "U2P, Document Collection, Presentation"
        prs.core_properties.category = "Document Collection"
        prs.core_properties.comments = f"Generated by U2P on {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}"
    
    others = []
    
    # Sort records by order
    sorted_records = sorted(records, key=lambda x: x.get("order", 0))
    
    # Add title slide if presentation title is provided
    if presentation_title:
        title_slide = prs.slides.add_slide(prs.slide_layouts[6])  # Use blank layout for full control
        
        # Calculate center position
        slide_width = prs.slide_width
        slide_height = prs.slide_height
        
        # Add title textbox centered
        title_box = title_slide.shapes.add_textbox(
            Inches(1.0), Inches(2.0),  # Left, Top position
            slide_width - Inches(2.0), Inches(1.5)  # Width, Height
        )
        title_frame = title_box.text_frame
        title_frame.clear()
        p = title_frame.paragraphs[0]
        p.text = presentation_title
        p.font.size = Pt(36)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
        
        # Add subtitle content below title
        subtitle_box = title_slide.shapes.add_textbox(
            Inches(1.0), Inches(3.8),  # Left, Top position
            slide_width - Inches(2.0), slide_height - Inches(5.0)  # Width, Height
        )
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.clear()
        
        # Add generation info
        p1 = subtitle_frame.paragraphs[0]
        p1.text = f"Generated by U2P - Upload to Present"
        p1.font.size = Pt(16)
        p1.font.color.rgb = RGBColor(128, 128, 128)
        p1.alignment = PP_ALIGN.CENTER
        
        # Add batch ID and date
        p2 = subtitle_frame.add_paragraph()
        p2.text = f"Batch ID: {batch_id} • {datetime.now().strftime('%B %d, %Y')}"
        p2.font.size = Pt(14)
        p2.font.color.rgb = RGBColor(128, 128, 128)
        p2.alignment = PP_ALIGN.CENTER
        
        # Add file count
        p3 = subtitle_frame.add_paragraph()
        p3.text = f"Total Files: {len(records)}"
        p3.font.size = Pt(14)
        p3.font.color.rgb = RGBColor(128, 128, 128)
        p3.alignment = PP_ALIGN.CENTER
        
        # Add file list
        if len(records) <= 8:  # Only show list if 8 or fewer files
            p4 = subtitle_frame.add_paragraph()
            p4.text = "Files Included:"
            p4.font.size = Pt(12)
            p4.font.color.rgb = RGBColor(100, 100, 100)
            p4.alignment = PP_ALIGN.CENTER
            
            # Add each file name with custom titles if provided
            for i, rec in enumerate(records[:6]):  # Limit to 6 files to avoid crowding
                p_file = subtitle_frame.add_paragraph()
                custom_title = rec.get("title", "")
                
                # Check if it's a camera image
                is_camera = rec.get("is_camera", False)
                camera_display = " (Camera)" if is_camera else ""
                
                if custom_title and custom_title != rec["name"]:
                    # Show custom title - filename - camera format
                    p_file.text = f"• {custom_title} - {rec['name']}{camera_display}"
                else:
                    # Show filename - camera format
                    p_file.text = f"• {rec['name']}{camera_display}"
                p_file.font.size = Pt(10)
                p_file.font.color.rgb = RGBColor(100, 100, 100)
                p_file.alignment = PP_ALIGN.CENTER
                
            if len(records) > 6:
                p_more = subtitle_frame.add_paragraph()
                p_more.text = f"... and {len(records) - 6} more files"
                p_more.font.size = Pt(10)
                p_more.font.color.rgb = RGBColor(100, 100, 100)
                p_more.alignment = PP_ALIGN.CENTER
        else:
            # For many files, just show count
            p4 = subtitle_frame.add_paragraph()
            p4.text = f"Files Included: {len(records)} documents"
            p4.font.size = Pt(12)
            p4.font.color.rgb = RGBColor(100, 100, 100)
            p4.alignment = PP_ALIGN.CENTER

    for rec in sorted_records:
        fname = rec["name"]
        data: bytes = rec["bytes"]
        mime = rec["mime"]

        if is_image(fname):
            # Always use blank layout to avoid placeholder text
            slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
            
            # Add title manually if enabled
            if show_titles:
                title_text = rec.get("title", fname)
                # Add title textbox manually
                title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), 
                                                   prs.slide_width - Inches(1.0), Inches(0.8))
                title_frame = title_box.text_frame
                title_frame.clear()
                p = title_frame.paragraphs[0]
                p.text = title_text
                p.font.size = Pt(18)
                p.font.bold = False
                p.alignment = PP_ALIGN.CENTER
                st.info(f"Added title: {title_text}")  # Debug message
            
            try:
                # Process image for PowerPoint compatibility
                processed_data = process_image_for_powerpoint(data, rec['name'])
                img = Image.open(io.BytesIO(processed_data))
                
                left, top, w, h = fit_image_on_slide(prs, img, has_title=show_titles)
                # python-pptx needs a file-like, but we can give it the bytes directly
                prs.slides[-1].shapes.add_picture(io.BytesIO(processed_data), left, top, width=w, height=h)
            except UnidentifiedImageError:
                others.append(rec)  # if it fails to parse, list as other
            except Exception as e:
                st.warning(f"⚠️ Could not process image '{rec['name']}': {str(e)}")
                others.append(rec)
            
            # Add footer to slide
            add_footer_to_slide(slide, batch_id, prs.slide_width, prs.slide_height)
        elif is_pdf(fname) and PDF_SUPPORT:
            try:
                # Convert PDF pages to images using PyMuPDF
                pdf_document = fitz.open(stream=data, filetype="pdf")
                for page_num in range(len(pdf_document)):
                    page = pdf_document.load_page(page_num)
                    # Render page to image with high DPI for quality
                    mat = fitz.Matrix(2.0, 2.0)  # 2x zoom for 150 DPI equivalent
                    pix = page.get_pixmap(matrix=mat)
                    # Convert to PIL Image
                    img_data = pix.tobytes("png")
                    pdf_img = Image.open(io.BytesIO(img_data))
                    
                    # Always use blank layout to avoid placeholder text
                    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
                    
                    # Add title manually if enabled (with page number for multi-page PDFs)
                    if show_titles:
                        base_title = rec.get("title", fname)
                        if len(pdf_document) > 1:
                            title_text = f"{base_title} (Page {page_num + 1})"
                        else:
                            title_text = base_title
                        
                        # Add title textbox manually
                        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), 
                                                           prs.slide_width - Inches(1.0), Inches(0.8))
                        title_frame = title_box.text_frame
                        title_frame.clear()
                        p = title_frame.paragraphs[0]
                        p.text = title_text
                        p.font.size = Pt(18)
                        p.font.bold = False
                        p.alignment = PP_ALIGN.CENTER
                        st.info(f"Added PDF title: {title_text}")  # Debug message
                    
                    left, top, w, h = fit_image_on_slide(prs, pdf_img, has_title=show_titles)
                    # Add image to slide
                    img_bytes = io.BytesIO()
                    pdf_img.save(img_bytes, format='PNG')
                    img_bytes.seek(0)
                    prs.slides[-1].shapes.add_picture(img_bytes, left, top, width=w, height=h)
                    
                    # Add footer to slide
                    add_footer_to_slide(slide, batch_id, prs.slide_width, prs.slide_height)
                
                pdf_document.close()
            except Exception as e:
                st.warning(f"Could not process PDF '{fname}': {str(e)}")
                others.append(rec)
        else:
            others.append(rec)

    if others:
        slide = prs.slides.add_slide(prs.slide_layouts[5])  # title only
        slide.shapes.title.text = "Attached Files"
        # Adjust textbox for portrait format
        txbox = slide.shapes.add_textbox(Inches(0.8), Inches(1.6),
                                         prs.slide_width - Inches(1.6),
                                         prs.slide_height - Inches(2.5))
        tf = txbox.text_frame
        tf.word_wrap = True
        tf.clear()

        # Add each filename as a bullet (no links—everything is in memory)
        for i, rec in enumerate(others):
            p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
            run = p.add_run()
            run.text = f"• {rec['name']}"
            run.font.size = Pt(16)
            p.alignment = PP_ALIGN.LEFT
            p.space_after = Pt(6)
        
        # Add footer to summary slide
        add_footer_to_slide(slide, batch_id, prs.slide_width, prs.slide_height)

    out = io.BytesIO()
    prs.save(out)
    out.seek(0)
    return out.read()

def build_zip_from_records(records: List[Dict], batch_id: str) -> bytes:
    """
    Creates an in-memory ZIP of the originals.
    """
    out = io.BytesIO()
    with zipfile.ZipFile(out, mode="w", compression=zipfile.ZIP_DEFLATED) as zf:
        for rec in records:
            zf.writestr(f"{batch_id}/{rec['name']}", rec["bytes"])
    out.seek(0)
    return out.read()

# ---------- UI ----------
st.set_page_config(page_title="U2P - Upload to Present", page_icon="📊", layout="centered")
st.title("📊 U2P - Upload to Present")
st.caption("Upload anything. Present everything. No disk, no cloud—just memory.")

# Session state
if "batch_id" not in st.session_state:
    st.session_state.batch_id = timestamp_slug()
if "records" not in st.session_state:
    st.session_state.records = []  # list of dicts: {name, bytes, mime, size, title, order}
if "show_titles" not in st.session_state:
    st.session_state.show_titles = False
if "current_step" not in st.session_state:
    st.session_state.current_step = 1
if "ppt_bytes" not in st.session_state:
    st.session_state.ppt_bytes = None
if "zip_bytes" not in st.session_state:
    st.session_state.zip_bytes = None
if "presentation_title" not in st.session_state:
    st.session_state.presentation_title = ""

# Wizard interface
st.info(f"Batch ID: {st.session_state.batch_id}")

# Step indicator
steps = ["📁 Upload Files", "⚙️ Organize & Configure", "📊 Download"]
current_step = st.session_state.current_step

# Progress bar
progress = current_step / len(steps)
st.progress(progress)

# Step navigation
col1, col2, col3 = st.columns(3)
with col1:
    if st.button("📁 Step 1: Upload", disabled=current_step==1):
        st.session_state.current_step = 1
        st.rerun()
with col2:
    if st.button("⚙️ Step 2: Organize", disabled=current_step==2 or len(st.session_state.records)==0):
        st.session_state.current_step = 2
        st.rerun()
with col3:
    if st.button("📊 Step 3: Download", disabled=current_step==3 or len(st.session_state.records)==0):
        st.session_state.current_step = 3
        st.rerun()

st.divider()

# Step 1: Upload Files
if current_step == 1:
    st.subheader("📁 Step 1: Upload Files")
    st.write(
        "• **Images** (PNG, JPG, GIF, etc.) become individual slides in the PowerPoint deck.\n"
        "• **PDFs** - each page becomes a separate slide in the PowerPoint deck.\n"
        "• **Videos, documents, and other files** are listed on a summary slide.\n"
        "• Nothing is written to disk or cloud—everything stays in memory."
    )
    
    uploaded = st.file_uploader(
        "Select files",
        accept_multiple_files=True,
        label_visibility="collapsed",
        help="Upload images, PDFs, videos, documents, or any other files."
    )
    
    col1, col2 = st.columns([1, 1])
    with col1:
        add_btn = st.button("Next: Organize Files", type="primary", disabled=not uploaded)
    with col2:
        reset_btn = st.button("Start new batch")
    
    if add_btn and uploaded:
        new_count = 0
        for f in uploaded:
            file_bytes = f.read()
            metadata = extract_file_metadata(file_bytes, f.name)
            
            rec = {
                "name": f.name,
                "mime": f.type or "application/octet-stream",
                "bytes": file_bytes,
                "size": f.size,
                "title": f.name.rsplit(".", 1)[0] if "." in f.name else f.name,
                "order": len(st.session_state.records) + new_count,
                "upload_time": metadata["upload_time"],
                "is_camera": metadata["is_camera"],
            }
            st.session_state.records.append(rec)
            new_count += 1
            time.sleep(0.01)
        st.success(f"Added {new_count} file(s) to batch.")
        # Automatically advance to next step
        st.session_state.current_step = 2
        st.rerun()
    
    if reset_btn:
        st.session_state.batch_id = timestamp_slug()
        st.session_state.records = []
        st.rerun()
    
    # Show current files
    if st.session_state.records:
        st.subheader("📋 Current Files")
        for i, r in enumerate(st.session_state.records):
            size_kb = f"{(r['size'] or len(r['bytes']))/1024:.1f} KB"
            
            # Check if it's a camera image
            is_camera = r.get("is_camera", False)
            camera_display = " (Camera)" if is_camera else ""
            
            st.write(f"✅ **{r['name']}{camera_display}** • {r['mime']} • {size_kb}")
        
        if st.button("Next: Organize Files", type="primary"):
            st.session_state.current_step = 2
            st.rerun()

# Step 2: Organize & Configure
elif current_step == 2:
    st.subheader("⚙️ Step 2: Organize & Configure")
    
    # Title toggle
    col1, col2 = st.columns([1, 1])
    with col1:
        st.session_state.show_titles = st.checkbox("Show titles on slides", value=st.session_state.show_titles)
    with col2:
        if st.button("Reset order"):
            for i, rec in enumerate(st.session_state.records):
                rec["order"] = i
            st.rerun()
    
    # Presentation title
    st.subheader("📋 Presentation Settings")
    st.session_state.presentation_title = st.text_input(
        "Presentation Title (optional)",
        value=st.session_state.presentation_title,
        placeholder="e.g., Document Collection, Project Review, Meeting Notes",
        help="This will be set as the PowerPoint file title and create a title slide"
    )
    
    # Sort records by order
    sorted_records = sorted(st.session_state.records, key=lambda x: x["order"])
    
    # File management interface
    for i, r in enumerate(sorted_records):
        with st.expander(f"📄 {r['name']} (Order: {i+1})", expanded=False):
            col1, col2, col3, col4 = st.columns([2, 1, 1, 1])
            
            with col1:
                new_title = st.text_input(
                    "Slide title:", 
                    value=r["title"], 
                    key=f"title_{i}",
                    help="Title to show on this slide (if enabled)"
                )
                r["title"] = new_title
                
                size_kb = f"{(r['size'] or len(r['bytes']))/1024:.1f} KB"
                
                # Check if it's a camera image
                is_camera = r.get("is_camera", False)
                camera_display = " (Camera)" if is_camera else ""
                
                st.caption(f"Type: {r['mime']} • Size: {size_kb}{camera_display}")
            
            with col2:
                if st.button("⬆️", key=f"up_{i}", disabled=i==0):
                    if i > 0:
                        sorted_records[i]["order"], sorted_records[i-1]["order"] = sorted_records[i-1]["order"], sorted_records[i]["order"]
                        st.rerun()
            
            with col3:
                if st.button("⬇️", key=f"down_{i}", disabled=i==len(sorted_records)-1):
                    if i < len(sorted_records) - 1:
                        sorted_records[i]["order"], sorted_records[i+1]["order"] = sorted_records[i+1]["order"], sorted_records[i]["order"]
                        st.rerun()
            
            with col4:
                if st.button("✕", key=f"remove_{i}", help="Remove this file from batch"):
                    st.session_state.records.remove(r)
                    st.rerun()
    
    # Quick order display
    st.caption("**Current order:** " + " → ".join([f"{i+1}. {r['name']}" for i, r in enumerate(sorted_records)]))
    
    col1, col2 = st.columns([1, 1])
    with col1:
        if st.button("← Back to Upload"):
            st.session_state.current_step = 1
            st.rerun()
    with col2:
        if st.button("Next: Download", type="primary"):
            st.session_state.current_step = 3
            st.rerun()

# Step 3: Download
elif current_step == 3:
    st.subheader("📊 Step 3: Download")
    
    # Summary
    st.write(f"**Ready to download presentation with {len(st.session_state.records)} files**")
    if st.session_state.show_titles:
        st.write("✅ Titles will be shown on slides")
    else:
        st.write("ℹ️ No titles will be shown on slides")
    
    # Download options
    col1, col2 = st.columns(2)
    
    with col1:
        st.subheader("📄 PowerPoint Presentation")
        # Use download_button directly as the main button
        ppt_bytes = build_pptx_from_records(st.session_state.records, st.session_state.show_titles, st.session_state.batch_id, st.session_state.presentation_title)
        # Generate filename based on presentation title
        if st.session_state.presentation_title:
            # Clean the title for filename use
            clean_title = "".join(c for c in st.session_state.presentation_title if c.isalnum() or c in (' ', '-', '_')).rstrip()
            clean_title = clean_title.replace(' ', '_')[:30]  # Limit length
            filename = f"{clean_title}_{st.session_state.batch_id}.pptx"
        else:
            filename = f"u2p_{st.session_state.batch_id}.pptx"
        
        st.download_button(
            "Download Presentation",
            data=ppt_bytes,
            file_name=filename,
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            key="ppt_download",
            use_container_width=True
        )
    
    with col2:
        st.subheader("📦 Original Files")
        # Use download_button directly as the main button
        zip_bytes = build_zip_from_records(st.session_state.records, st.session_state.batch_id)
        st.download_button(
            "Download ZIP",
            data=zip_bytes,
            file_name=f"u2p_{st.session_state.batch_id}.zip",
            mime="application/zip",
            key="zip_download",
            use_container_width=True
        )
    
    # Navigation
    col1, col2 = st.columns([1, 1])
    with col1:
        if st.button("← Back to Organize"):
            st.session_state.current_step = 2
            st.rerun()
    with col2:
        if st.button("Start New Batch", type="secondary"):
            st.session_state.batch_id = timestamp_slug()
            st.session_state.records = []
            st.session_state.current_step = 1
            st.rerun()

st.divider()
st.caption(
    "Note: Streamlit Cloud has memory limits. For very large batches, consider paging or client-side pruning."
)
